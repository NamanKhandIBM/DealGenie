"""
DealGenie Pitch Deck builder
Produces: DealGenie_Deck.pptx  (4 slides)

Slide 1 — Title
Slide 2 — Problem / Solution  (3 × 2 grid with icons)
Slide 3 — Demo  (what to show, flow steps)
Slide 4 — Feedback / Next Steps

Run:  python3 scripts/build_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy, os

# ── Palette ──────────────────────────────────────────────────────────────────
IBM_BLUE       = RGBColor(0x0F, 0x62, 0xFE)
IBM_DARK       = RGBColor(0x16, 0x16, 0x16)
IBM_MID        = RGBColor(0x52, 0x52, 0x52)
IBM_LIGHT_BG   = RGBColor(0xF4, 0xF7, 0xFF)
IBM_RULE       = RGBColor(0xD0, 0xD0, 0xD0)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
RED            = RGBColor(0xDA, 0x1E, 0x28)
GREEN          = RGBColor(0x24, 0xA1, 0x48)
AMBER          = RGBColor(0xF1, 0xC2, 0x1B)
PROB_BG        = RGBColor(0xFF, 0xF5, 0xF5)
SOLN_BG        = RGBColor(0xF0, 0xF6, 0xFF)

# ── Slide dimensions  (16:9 widescreen) ──────────────────────────────────────
W = Inches(13.33)
H = Inches(7.50)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank


# ══════════════════════════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def add_rect(slide, left, top, width, height,
             fill=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.width = line_width
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width if line_width else Pt(0.75)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, text, left, top, width, height,
             size=Pt(12), bold=False, color=IBM_DARK,
             align=PP_ALIGN.LEFT, wrap=True,
             italic=False, v_anchor="top"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    from pptx.enum.text import MSO_ANCHOR
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "IBM Plex Sans"
    return txBox


def add_line(slide, x1, y1, x2, y2, color=IBM_RULE, width=Pt(0.75)):
    """Thin divider line."""
    from pptx.util import Emu
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = width
    return conn


def footer(slide, page_num):
    """Adds a thin bottom bar with page number."""
    bar_h = Inches(0.28)
    add_rect(slide, 0, H - bar_h, W, bar_h, fill=WHITE,
             line_color=IBM_RULE, line_width=Pt(0.5))
    add_text(slide, "IBM INTERNAL — CONFIDENTIAL",
             Inches(0.4), H - bar_h + Inches(0.04),
             Inches(4), bar_h,
             size=Pt(7), color=RGBColor(0xA8, 0xA8, 0xA8))
    add_text(slide, str(page_num),
             W - Inches(0.6), H - bar_h + Inches(0.04),
             Inches(0.4), bar_h,
             size=Pt(8), bold=True, color=IBM_MID,
             align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════

s1 = prs.slides.add_slide(blank_layout)

# Full dark background
add_rect(s1, 0, 0, W, H, fill=IBM_DARK)

# Blue accent bar left
add_rect(s1, 0, 0, Inches(0.18), H, fill=IBM_BLUE)

# Subtle grid pattern (light rectangles)
for i in range(6):
    add_rect(s1,
             Inches(1.8 + i * 1.9), Inches(0.5),
             Inches(1.5), H - Inches(1.0),
             fill=RGBColor(0x20, 0x20, 0x20),
             line_color=RGBColor(0x2A, 0x2A, 0x2A),
             line_width=Pt(0.25))

# IBM Blue glow block behind title
add_rect(s1, Inches(0.55), Inches(2.1), Inches(7.5), Inches(0.08),
         fill=IBM_BLUE)

# Product name badge
badge = add_rect(s1, Inches(0.55), Inches(1.55), Inches(1.9), Inches(0.36),
                 fill=IBM_BLUE)

add_text(s1, "IBM INTERNAL", Inches(0.6), Inches(1.58),
         Inches(1.8), Inches(0.30),
         size=Pt(8), bold=True, color=WHITE,
         align=PP_ALIGN.LEFT)

# Main title
add_text(s1, "Deal Genie",
         Inches(0.55), Inches(2.2), Inches(9), Inches(1.4),
         size=Pt(66), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Subtitle
add_text(s1,
         "AI-powered quoting assistant for IBM Security sellers.\n"
         "Discovery → Part Numbers → CPQ in under 2 minutes.",
         Inches(0.55), Inches(3.65), Inches(8), Inches(1.0),
         size=Pt(17), color=RGBColor(0xC6, 0xC6, 0xC6), align=PP_ALIGN.LEFT)

# Powered-by row
add_text(s1, "Powered by",
         Inches(0.55), Inches(5.0), Inches(1.5), Inches(0.4),
         size=Pt(10), color=RGBColor(0x78, 0x78, 0x78))
add_text(s1, "watsonx.ai",
         Inches(1.45), Inches(5.0), Inches(2), Inches(0.4),
         size=Pt(10), bold=True, color=IBM_BLUE)
add_text(s1, "  ·  IBM Cloudant  ·  IBM Marketplace API",
         Inches(2.7), Inches(5.0), Inches(5), Inches(0.4),
         size=Pt(10), color=RGBColor(0x78, 0x78, 0x78))

# Products covered
for i, (name, col) in enumerate([
    ("IBM Security Verify", RGBColor(0x4D, 0x8E, 0xF8)),
    ("NS1 Connect",         RGBColor(0x10, 0xB9, 0x81)),
    ("IBM HashiCorp Vault", RGBColor(0xA8, 0x55, 0xF7)),
]):
    bx = Inches(0.55 + i * 2.55)
    add_rect(s1, bx, Inches(5.7), Inches(2.3), Inches(0.46),
             fill=RGBColor(0x22, 0x22, 0x22),
             line_color=col, line_width=Pt(1))
    add_text(s1, name, bx + Inches(0.12), Inches(5.75),
             Inches(2.1), Inches(0.36),
             size=Pt(10), bold=True, color=col)

# 50 tests badge
add_rect(s1, Inches(8.4), Inches(5.7), Inches(2.2), Inches(0.46),
         fill=RGBColor(0x22, 0x22, 0x22),
         line_color=GREEN, line_width=Pt(1))
add_text(s1, "50 tests · 0 failures",
         Inches(8.52), Inches(5.75), Inches(2.1), Inches(0.36),
         size=Pt(10), bold=True, color=GREEN)

footer(s1, 1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM / SOLUTION
# ══════════════════════════════════════════════════════════════════════════════

s2 = prs.slides.add_slide(blank_layout)
add_rect(s2, 0, 0, W, H, fill=WHITE)

# ── Left headline panel ───────────────────────────────────────────────────────
panel_w = Inches(2.6)
add_rect(s2, 0, 0, panel_w, H - Inches(0.28), fill=WHITE,
         line_color=IBM_RULE, line_width=Pt(0.5))
add_rect(s2, 0, 0, Inches(0.12), H - Inches(0.28), fill=IBM_BLUE)

add_text(s2, "Quoting IBM products is harder than it should be",
         Inches(0.3), Inches(0.55), Inches(2.1), Inches(2.0),
         size=Pt(22), bold=True, color=IBM_DARK)

add_text(s2, "Deal Genie solves it",
         Inches(0.3), H - Inches(0.85), Inches(2.1), Inches(0.4),
         size=Pt(9), bold=True,
         color=RGBColor(0x6F, 0x6F, 0x6F))

# ── Grid setup ────────────────────────────────────────────────────────────────
gx     = panel_w                           # grid starts here
gw     = W - panel_w                       # total grid width
gh     = H - Inches(0.28)                  # grid height (minus footer)
col_w  = gw / 3
row_h  = gh / 2

# Column dividers
for i in (1, 2):
    add_line(s2, gx + col_w * i, 0, gx + col_w * i, gh,
             color=IBM_RULE, width=Pt(0.5))

# Row divider
add_line(s2, gx, row_h, W, row_h, color=IBM_RULE, width=Pt(0.5))

# Row background tints
add_rect(s2, gx, 0, gw, row_h,
         fill=RGBColor(0xFF, 0xFB, 0xFB))          # subtle warm red for problems
add_rect(s2, gx, row_h, gw, row_h,
         fill=RGBColor(0xF4, 0xF8, 0xFF))          # subtle blue for solutions

# ── Row header strips ─────────────────────────────────────────────────────────
prob_strip_h = Inches(0.26)
soln_strip_h = Inches(0.26)
add_rect(s2, gx, 0, gw, prob_strip_h, fill=RED)
add_text(s2, "THE PROBLEMS",
         gx + Inches(0.25), Inches(0.05), gw - Inches(0.5), prob_strip_h,
         size=Pt(8), bold=True, color=WHITE,
         align=PP_ALIGN.LEFT)

add_rect(s2, gx, row_h, gw, soln_strip_h, fill=IBM_BLUE)
add_text(s2, "HOW DEAL GENIE SOLVES THEM",
         gx + Inches(0.25), row_h + Inches(0.05),
         gw - Inches(0.5), soln_strip_h,
         size=Pt(8), bold=True, color=WHITE,
         align=PP_ALIGN.LEFT)

# ── Cell content helper ───────────────────────────────────────────────────────
PROBLEMS = [
    (
        "No guided discovery",
        "Sellers didn't know what to ask clients.\n"
        "Every new seller had to call a specialist just to start a quote — "
        "blocking deals and slowing ramp time.",
        "P"
    ),
    (
        "Quotes lost — work repeated every call",
        "Nothing saved between sessions.\n"
        "Work was recreated from scratch on every discovery call, "
        "burning time and creating inconsistency across the team.",
        "P"
    ),
    (
        "CPQ rejections from wrong SKUs",
        "One missing required part or wrong tier rejected the whole order.\n"
        "Approval round-trips added days to deal cycles and eroded seller confidence.",
        "P"
    ),
]

SOLUTIONS = [
    (
        "Guided discovery built in",
        "Smart, context-aware questions walk any seller through every product requirement. "
        "The right SKUs are produced automatically — no specialist, no guesswork, no delay.",
        "S"
    ),
    (
        "Full history saved to IBM Cloudant",
        "Every quote persisted with full conversation context. "
        "Reload any past deal in one click — return to follow-up calls fully prepared.",
        "S"
    ),
    (
        "Exact SKUs, confirmed from IBM APIs",
        "Prices from IBM Marketplace API. Required parts auto-inserted. "
        "CPQ-ready CSV export — paste straight in, zero transcription errors, zero rejections.",
        "S"
    ),
]

ICONS_P = ["⚠", "⏳", "✗"]    # rendered as text; replaced by shape below
ICONS_S = ["✦", "☁", "✓"]

def cell_content(slide, row, col, title, body, kind):
    """Fill one grid cell with number tag, title and body."""
    cx = gx + col * col_w
    cy = (0 if row == 0 else row_h)
    strip_h = prob_strip_h if row == 0 else soln_strip_h

    pad_top = strip_h + Inches(0.18)

    # Number badge
    num = col + 1
    tag_color = RED if kind == "P" else IBM_BLUE
    tag_bg    = RGBColor(0xFF, 0xE5, 0xE5) if kind == "P" else RGBColor(0xD4, 0xE6, 0xFF)

    badge_w, badge_h = Inches(0.52), Inches(0.24)
    add_rect(slide, cx + Inches(0.22), cy + pad_top,
             badge_w, badge_h,
             fill=tag_bg, line_color=tag_color, line_width=Pt(0.6))
    label = f"{'Problem' if kind == 'P' else 'Solution'} 0{num}"
    add_text(slide, label,
             cx + Inches(0.23), cy + pad_top + Inches(0.02),
             badge_w - Inches(0.04), badge_h,
             size=Pt(6.5), bold=True, color=tag_color)

    # Icon circle
    icon_top = cy + pad_top + badge_h + Inches(0.12)
    ic_col = RED if kind == "P" else IBM_BLUE
    ic_bg  = RGBColor(0xFF, 0xEB, 0xEB) if kind == "P" else RGBColor(0xEB, 0xF3, 0xFF)
    add_rect(slide, cx + Inches(0.22), icon_top,
             Inches(0.38), Inches(0.38),
             fill=ic_bg, line_color=ic_col, line_width=Pt(0.75))

    # Title
    title_top = icon_top + Inches(0.48)
    add_text(slide, title,
             cx + Inches(0.22), title_top,
             col_w - Inches(0.44), Inches(0.62),
             size=Pt(13), bold=True, color=IBM_DARK)

    # Body
    body_top = title_top + Inches(0.68)
    add_text(slide, body,
             cx + Inches(0.22), body_top,
             col_w - Inches(0.44), row_h - body_top + cy - Inches(0.2),
             size=Pt(10.5), color=IBM_MID)


for col, (title, body, kind) in enumerate(PROBLEMS):
    cell_content(s2, 0, col, title, body, kind)

for col, (title, body, kind) in enumerate(SOLUTIONS):
    cell_content(s2, 1, col, title, body, kind)

# Arrow pills between rows
for col in range(3):
    ax = gx + col * col_w + col_w / 2 - Inches(0.18)
    ay = row_h - Inches(0.18)
    add_rect(s2, ax, ay, Inches(0.36), Inches(0.36),
             fill=IBM_BLUE, line_color=WHITE, line_width=Pt(2))
    add_text(s2, "↓", ax, ay - Inches(0.02), Inches(0.36), Inches(0.38),
             size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

footer(s2, 2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — LIVE DEMO
# ══════════════════════════════════════════════════════════════════════════════

s3 = prs.slides.add_slide(blank_layout)
add_rect(s3, 0, 0, W, H, fill=RGBColor(0xF7, 0xF8, 0xFA))

# Top bar
add_rect(s3, 0, 0, W, Inches(0.9), fill=IBM_DARK)
add_rect(s3, 0, 0, Inches(0.12), Inches(0.9), fill=IBM_BLUE)
add_text(s3, "Live Demo", Inches(0.3), Inches(0.12),
         Inches(6), Inches(0.7), size=Pt(28), bold=True, color=WHITE)
add_text(s3, "Requirements → Part Numbers → CPQ in under 2 minutes",
         Inches(6.2), Inches(0.28), Inches(6.8), Inches(0.5),
         size=Pt(13), color=RGBColor(0xA8, 0xC8, 0xFF),
         align=PP_ALIGN.RIGHT)

DEMO_STEPS = [
    ("01", "Log In",
     "Navigate to localhost:3000\nPassword: garland\nPassword-gated — IBM sellers only",
     IBM_BLUE),
    ("02", "Select IBM Security Verify",
     "Choose Start Quoting\n6 questions → MAU calculated automatically\nPart numbers + list price in 90 seconds",
     IBM_BLUE),
    ("03", "Export CSV + Save Quote",
     "Download CPQ-ready file — paste directly in\nSave to IBM Cloudant\nFull context persisted for follow-up calls",
     IBM_BLUE),
    ("04", "NS1 Connect — Premium Quote",
     "2,000 MQ scenario → routes to Premium tier\nTier-aware questions only (GSLB, DDoS appear at 1B+ MQ)\nAuto-inserts required SLA part D0GNDZX",
     IBM_BLUE),
    ("05", "Compare Scenarios  ← KEY",
     "Toggle GSLB / DDoS on/off live\nRunning total updates instantly\nDeterministic — zero AI — safe to use on a client call",
     AMBER),
    ("06", "Best Practices AI + Client Mode",
     "watsonx.ai SME answers product questions\nLive IBM Docs context injected\nClient Mode: AI speaks directly to the client",
     IBM_BLUE),
]

cols = 3
rows = 2
card_w = (W - Inches(0.5)) / cols - Inches(0.15)
card_h = (H - Inches(1.2)) / rows - Inches(0.12)

for i, (num, title, body, accent) in enumerate(DEMO_STEPS):
    col = i % cols
    row = i // cols
    cx = Inches(0.25) + col * (card_w + Inches(0.15))
    cy = Inches(1.0) + row * (card_h + Inches(0.12))

    # Card background
    add_rect(s3, cx, cy, card_w, card_h, fill=WHITE,
             line_color=IBM_RULE, line_width=Pt(0.5))
    # Accent left stripe
    add_rect(s3, cx, cy, Inches(0.06), card_h, fill=accent)
    # Number
    add_text(s3, num, cx + Inches(0.16), cy + Inches(0.14),
             Inches(0.5), Inches(0.38),
             size=Pt(22), bold=True, color=accent)
    # Title
    add_text(s3, title, cx + Inches(0.16), cy + Inches(0.52),
             card_w - Inches(0.22), Inches(0.38),
             size=Pt(13.5), bold=True, color=IBM_DARK)
    # Divider
    add_line(s3, cx + Inches(0.16), cy + Inches(0.94),
             cx + card_w - Inches(0.1), cy + Inches(0.94),
             color=IBM_RULE, width=Pt(0.4))
    # Body
    add_text(s3, body, cx + Inches(0.16), cy + Inches(1.0),
             card_w - Inches(0.22), card_h - Inches(1.1),
             size=Pt(10), color=IBM_MID)

# Bottom note
add_rect(s3, Inches(0.25), H - Inches(0.55), W - Inches(0.5), Inches(0.26),
         fill=RGBColor(0xEB, 0xF3, 0xFF),
         line_color=IBM_BLUE, line_width=Pt(0.6))
add_text(s3,
         "All prices are LIST — confirmed from IBM Marketplace API.  "
         "Deterministic engine · 50 automated tests · Zero AI in pricing calculations.",
         Inches(0.35), H - Inches(0.52), W - Inches(0.7), Inches(0.24),
         size=Pt(8.5), color=IBM_BLUE)

footer(s3, 3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — FEEDBACK / NEXT STEPS
# ══════════════════════════════════════════════════════════════════════════════

s4 = prs.slides.add_slide(blank_layout)
add_rect(s4, 0, 0, W, H, fill=WHITE)

# Top bar
add_rect(s4, 0, 0, W, Inches(0.9), fill=IBM_DARK)
add_rect(s4, 0, 0, Inches(0.12), Inches(0.9), fill=IBM_BLUE)
add_text(s4, "Feedback & Next Steps", Inches(0.3), Inches(0.12),
         Inches(9), Inches(0.7), size=Pt(28), bold=True, color=WHITE)

# ── Two column layout ──
left_x  = Inches(0.45)
right_x = Inches(6.9)
col_w2  = Inches(6.0)
top_y   = Inches(1.1)

# LEFT — What we're asking for
add_rect(s4, left_x, top_y, col_w2, Inches(0.34), fill=IBM_BLUE)
add_text(s4, "What we're asking for",
         left_x + Inches(0.15), top_y + Inches(0.05),
         col_w2 - Inches(0.3), Inches(0.26),
         size=Pt(11), bold=True, color=WHITE)

ASKS = [
    ("Green light to deploy to the field team",
     "The tool is production-ready today. Let real sellers use it on real deals "
     "so we can gather feedback before scoping the next phase."),
    ("Input on hosting",
     "Containerised with Docker — ready for IBM Cloud Code Engine, IKS, or any "
     "internal IBM server in under an hour. Where should this live?"),
    ("A short pilot group",
     "3–5 sellers across NS1 and Verify. Two weeks of real deal usage "
     "will give us more signal than months of internal testing."),
]

for i, (ask_title, ask_body) in enumerate(ASKS):
    ay = top_y + Inches(0.45) + i * Inches(1.45)
    add_rect(s4, left_x, ay, col_w2, Inches(1.35),
             fill=RGBColor(0xF4, 0xF8, 0xFF),
             line_color=IBM_BLUE, line_width=Pt(0.6))
    add_rect(s4, left_x, ay, Inches(0.06), Inches(1.35), fill=IBM_BLUE)
    add_text(s4, ask_title,
             left_x + Inches(0.16), ay + Inches(0.1),
             col_w2 - Inches(0.25), Inches(0.35),
             size=Pt(12), bold=True, color=IBM_DARK)
    add_text(s4, ask_body,
             left_x + Inches(0.16), ay + Inches(0.45),
             col_w2 - Inches(0.25), Inches(0.82),
             size=Pt(10), color=IBM_MID)

# RIGHT — Roadmap
add_rect(s4, right_x, top_y, col_w2, Inches(0.34), fill=IBM_BLUE)
add_text(s4, "Roadmap — what's ready to build",
         right_x + Inches(0.15), top_y + Inches(0.05),
         col_w2 - Inches(0.3), Inches(0.26),
         size=Pt(11), bold=True, color=WHITE)

ROADMAP = [
    ("PDF Quote Export",      "~½ day",  GREEN,  "Client-facing quote document — Download PDF button"),
    ("Multi-product sessions", "~1 day",  GREEN,  "Quote Verify + NS1 + Vault in a single conversation"),
    ("Discount guardrails",   "~½ day",  GREEN,  "Warn seller when proposed discount exceeds IBM policy"),
    ("Salesforce CRM pre-fill","~1 day", AMBER,  "Pull opportunity data to pre-populate discovery · needs SF credentials"),
    ("CPQ direct push",       "TBD",     AMBER,  "Send finished quote into IBM CPQ · needs CPQ API access"),
]

for i, (rm_title, effort, col, desc) in enumerate(ROADMAP):
    ry = top_y + Inches(0.45) + i * Inches(1.02)
    add_rect(s4, right_x, ry, col_w2, Inches(0.92),
             fill=RGBColor(0xF7, 0xF8, 0xFA),
             line_color=IBM_RULE, line_width=Pt(0.5))
    add_rect(s4, right_x, ry, Inches(0.06), Inches(0.92), fill=col)
    add_text(s4, rm_title,
             right_x + Inches(0.16), ry + Inches(0.08),
             col_w2 - Inches(0.6), Inches(0.32),
             size=Pt(12), bold=True, color=IBM_DARK)
    # Effort badge
    add_rect(s4, right_x + col_w2 - Inches(0.8), ry + Inches(0.1),
             Inches(0.7), Inches(0.22),
             fill=col, line_color=col, line_width=Pt(0))
    add_text(s4, effort,
             right_x + col_w2 - Inches(0.78), ry + Inches(0.11),
             Inches(0.68), Inches(0.2),
             size=Pt(7.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s4, desc,
             right_x + Inches(0.16), ry + Inches(0.42),
             col_w2 - Inches(0.25), Inches(0.44),
             size=Pt(9.5), color=IBM_MID)

# Bottom closing line
add_rect(s4, Inches(0.45), H - Inches(0.58),
         W - Inches(0.9), Inches(0.28),
         fill=RGBColor(0x16, 0x16, 0x16))
add_text(s4,
         '"Right now, every NS1 and Verify quote requires an SME meeting and days of delay. '
         'DealGenie removes that friction entirely."',
         Inches(0.6), H - Inches(0.54),
         W - Inches(1.2), Inches(0.26),
         size=Pt(9), italic=True, color=RGBColor(0xC6, 0xC6, 0xC6))

footer(s4, 4)


# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════

out_path = os.path.join(os.path.dirname(__file__), "..", "DealGenie_Deck.pptx")
out_path = os.path.normpath(out_path)
prs.save(out_path)
print(f"Saved: {out_path}")
